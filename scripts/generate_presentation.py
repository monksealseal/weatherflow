#!/usr/bin/env python3
"""Generate WeatherFlow onboarding presentation as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──────────────────────────────────────────────────────────────────

BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD      = RGBColor(0x1A, 0x25, 0x3C)   # card surface
ACCENT       = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
ACCENT2      = RGBColor(0x4A, 0xDE, 0x80)   # green
ACCENT3      = RGBColor(0xFB, 0xBF, 0x24)   # amber
ACCENT4      = RGBColor(0xF8, 0x71, 0x71)   # coral
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY     = RGBColor(0x94, 0xA3, 0xB8)

WIDTH  = Inches(13.333)
HEIGHT = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, *,
                 font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, color=WHITE, bold=False, space_before=Pt(6),
             alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_card(slide, left, top, width, height, title, bullets, accent_color=ACCENT):
    """Add a card-style box with title and bullet points."""
    rect = add_rounded_rect(slide, left, top, width, height, BG_CARD)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    tf = add_text_box(slide, left + Inches(0.3), top + Inches(0.15),
                      width - Inches(0.5), Inches(0.4),
                      title, font_size=16, bold=True, color=accent_color)

    # Bullets
    for b in bullets:
        add_para(tf, b, font_size=13, color=LIGHT_GRAY, space_before=Pt(4))

    return rect


# ── Build Presentation ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = WIDTH
prs.slide_height = HEIGHT

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             "WeatherFlow", font_size=60, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.6),
             "Interactive Weather AI Platform",
             font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(0.6),
             "Train AI weather models  |  Real scientific data  |  No code required",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Badge row
badges = ["Python 3.8+", "PyTorch 2.0+", "Streamlit", "MIT License"]
badge_colors = [ACCENT, ACCENT3, ACCENT4, ACCENT2]
x_start = Inches(3.2)
for i, (label, c) in enumerate(zip(badges, badge_colors)):
    x = x_start + Inches(i * 1.9)
    r = add_rounded_rect(slide, x, Inches(5.0), Inches(1.6), Inches(0.4), BG_CARD)
    add_text_box(slide, x, Inches(5.02), Inches(1.6), Inches(0.36),
                 label, font_size=12, bold=True, color=c, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
             "v0.4.3  |  github.com/monksealseal/weatherflow",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is WeatherFlow?
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "What is WeatherFlow?", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "An interactive platform for training and experimenting with AI weather prediction models.",
             font_size=18, color=LIGHT_GRAY)

# 4 value cards
cards = [
    ("Train Real Neural Networks",
     ["Real PyTorch models on NCEP/ERA5 data",
      "Flow matching + physics-informed losses",
      "Watch training progress live"],
     ACCENT),
    ("Visualize Weather Predictions",
     ["Publication-quality weather maps",
      "Multi-day forecast animations",
      "SkewT diagrams, Hovmoller plots"],
     ACCENT2),
    ("Benchmark Against the Best",
     ["Compare to GraphCast, Pangu-Weather",
      "WeatherBench2 standard metrics",
      "Z500 RMSE, T850, wind scores"],
     ACCENT3),
    ("Learn Atmospheric Science",
     ["Graduate-level physics labs",
      "Interactive GCM simulations",
      "Rossby waves, vorticity, geostrophic flow"],
     ACCENT4),
]
for i, (title, bullets, color) in enumerate(cards):
    col = i % 4
    x = Inches(0.8) + col * Inches(3.1)
    add_card(slide, x, Inches(2.0), Inches(2.85), Inches(2.6), title, bullets, color)

add_text_box(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.8),
             'Key insight: The Streamlit app is not a demo with fake data. '
             'It executes the actual Python code from weatherflow/, applications/, and gcm/ modules. '
             'Every calculation, every model, every visualization runs real code.',
             font_size=16, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Quick Start
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Quick Start (10 seconds)", font_size=36, bold=True, color=WHITE)

# Option 1 card
add_card(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(3.0),
         "Option 1: Run Locally (Recommended)",
         ["git clone github.com/monksealseal/weatherflow.git",
          "cd weatherflow",
          "pip install -r streamlit_app/requirements.txt",
          "streamlit run streamlit_app/Home.py",
          "",
          "Open http://localhost:8501"],
         ACCENT)

# Option 2 card
add_card(slide, Inches(6.8), Inches(1.4), Inches(5.5), Inches(3.0),
         "Option 2: Streamlit Cloud",
         ["1. Fork the repository on GitHub",
          "2. Go to streamlit.io/cloud",
          "3. Connect your fork",
          "4. Set main file: streamlit_app/Home.py",
          "5. Deploy!"],
         ACCENT2)

# Dev tools card
add_card(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.2),
         "Developer Setup (full build system)",
         ["make install-all          Install Python + frontend + docs dependencies",
          "make dev-backend          Start FastAPI backend on port 8000",
          "make dev-frontend         Start React frontend on port 5173",
          "make test                 Run all 100+ tests",
          "make lint                 Run flake8 / black / isort / mypy",
          "make build                Build Python wheel + frontend + docs"],
         ACCENT3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Architecture Overview", font_size=36, bold=True, color=WHITE)

# Python core
add_card(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(2.5),
         "weatherflow/  (Python Core)",
         ["models/ - Flow matching, physics-guided nets",
          "training/ - Trainers, metrics, losses",
          "data/ - ERA5, NCEP data loaders",
          "physics/ - Atmospheric constraints",
          "inference/ - Prediction pipeline",
          "server/ - FastAPI REST API"],
         ACCENT)

# Applications
add_card(slide, Inches(4.6), Inches(1.4), Inches(3.8), Inches(2.5),
         "applications/  (Domain Apps)",
         ["renewable_energy/ - Wind & solar power",
          "extreme_event_analysis/ - Heatwaves, ARs",
          "educational/ - Teaching materials",
          "",
          "gcm/ - General Circulation Model",
          "Full atmospheric physics simulation"],
         ACCENT2)

# Frontend + Streamlit
add_card(slide, Inches(8.7), Inches(1.4), Inches(4.1), Inches(2.5),
         "User Interfaces",
         ["streamlit_app/ - 28 interactive pages",
          "   No-code access to all Python modules",
          "frontend/ - React + TypeScript SPA",
          "   Vite 5, Plotly.js, Three.js",
          "docs/ - MkDocs Material documentation",
          "   API reference + tutorials"],
         ACCENT3)

# Infrastructure
add_card(slide, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.8),
         "Infrastructure & CI/CD",
         ["Makefile - 40+ targets: unified entry point for everything",
          "GitHub Actions - tests, lint, build, deploy workflows",
          "Docker - docker-compose with demo/web/shell profiles",
          "PyPI - pip install weatherflow",
          "GitHub Pages - frontend deployment",
          "Pre-commit hooks - black, isort, flake8, mypy"],
         ACCENT4)

# Key tech
add_card(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.8),
         "Key Technologies",
         ["PyTorch 2.0+ - deep learning framework",
          "torchdiffeq - ODE solvers for flow matching",
          "xarray + zarr - scientific data handling",
          "FastAPI - async REST backend",
          "Streamlit - interactive web dashboards",
          "Plotly / Three.js - visualizations"],
         MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Streamlit App Pages (Part 1 - Core Workflows)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Streamlit App: Core Workflow Pages", font_size=36, bold=True, color=WHITE)

pages_core = [
    ("Home Dashboard", "Central hub: data status, model status, quick actions, benchmark leaderboard",
     ACCENT),
    ("Data Manager", "Load NCEP/ERA5 data with one click (7 MB sample), upload custom datasets, visualize in real-time",
     ACCENT),
    ("Training Workflow", "Configure architecture (layers, attention, physics), set hyperparameters, watch live loss curves",
     ACCENT2),
    ("Weather Prediction", "Load checkpoints, initialize from current conditions, generate multi-day forecasts",
     ACCENT2),
    ("Model Comparison", "Benchmark vs GraphCast / Pangu-Weather / FourCastNet using WeatherBench2 metrics",
     ACCENT3),
    ("Flow Matching", "Build WeatherFlowMatch models, visualize flow fields, interactive ODE integration demos",
     ACCENT3),
]
for i, (title, desc, color) in enumerate(pages_core):
    row = i // 3
    col = i % 3
    x = Inches(0.5) + col * Inches(4.15)
    y = Inches(1.4) + row * Inches(2.7)
    add_card(slide, x, y, Inches(3.9), Inches(2.3), title, [desc], color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Streamlit App Pages (Part 2 - Applications & Science)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Streamlit App: Applications & Science Pages", font_size=36, bold=True, color=WHITE)

pages_apps = [
    ("Wind Power", "Real turbine models (Vestas V90, Siemens SG 14), power curves, capacity factors, AEP estimates",
     ACCENT),
    ("Solar Power", "PV system configuration, irradiance calculations, temperature effects, grid integration",
     ACCENT),
    ("Extreme Events", "Heatwave detection, atmospheric river identification, extreme precipitation analysis",
     ACCENT4),
    ("GCM Simulation", "Full atmospheric physics: radiation, convection, boundary layers, Coriolis force",
     ACCENT2),
    ("Education", "Graduate-level labs: balanced flow, Rossby waves, vorticity, thermal wind, hydrostatic balance",
     ACCENT2),
    ("Hurricane Tracking", "IBTrACS / HURDAT2 data, satellite imagery, AI-powered wind field + intensity inference",
     ACCENT3),
]
for i, (title, desc, color) in enumerate(pages_apps):
    row = i // 3
    col = i % 3
    x = Inches(0.5) + col * Inches(4.15)
    y = Inches(1.4) + row * Inches(2.7)
    add_card(slide, x, y, Inches(3.9), Inches(2.3), title, [desc], color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — All 28 Streamlit Pages at a Glance
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "All 28 Streamlit Pages", font_size=36, bold=True, color=WHITE)

all_pages = [
    ("Home.py", "Dashboard & status"),
    ("0_Data_Manager", "Load weather data"),
    ("01_Live_Dashboard", "Real-time monitoring"),
    ("02_Research_Feed", "Latest research"),
    ("03_Training_Workflow", "Train AI models"),
    ("04_Visualization_Studio", "Create visualizations"),
    ("1_Wind_Power", "Wind energy forecasts"),
    ("2_Solar_Power", "Solar energy forecasts"),
    ("3_Extreme_Events", "Event detection"),
    ("4_Flow_Matching", "Flow model demos"),
    ("5_GCM_Simulation", "Climate modeling"),
    ("6_Education", "Physics labs"),
    ("7_Visualization", "Weather maps"),
    ("8_Physics_Losses", "Conservation laws"),
    ("9_Experiments", "Experiment tracking"),
    ("10_Model_Library", "Pre-trained models"),
    ("11_Training_Hub", "Advanced training"),
    ("12_Model_Comparison", "Benchmarking"),
    ("13_Publication_Vis", "Pub-quality plots"),
    ("14_GAIA_Functions", "Foundation model"),
    ("15_Research_Workbench", "Research tools"),
    ("16_WeatherBench2", "Standard metrics"),
    ("17_Weather_Prediction", "Generate forecasts"),
    ("18_Enterprise_Models", "Custom industry models"),
    ("20_Worldsphere", "AI command center"),
    ("21_Results_Gallery", "Results browser"),
    ("22_Hurricane_Tracking", "Hurricane analysis"),
    ("23_Tropic_World", "3D tropical sim"),
]

for i, (name, desc) in enumerate(all_pages):
    col = i // 10
    row = i % 10
    x = Inches(0.5) + col * Inches(4.3)
    y = Inches(1.3) + row * Inches(0.55)
    color = [ACCENT, ACCENT2, ACCENT3][col]
    tf = add_text_box(slide, x, y, Inches(4.1), Inches(0.5),
                      name, font_size=12, bold=True, color=color)
    add_para(tf, f"  {desc}", font_size=11, color=LIGHT_GRAY, space_before=Pt(0))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — How Flow Matching Works
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Core ML: How Flow Matching Works", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "WeatherFlow uses continuous normalizing flows (CNFs) to learn the evolution of weather states.",
             font_size=18, color=LIGHT_GRAY)

steps = [
    ("1. Data", "Load real atmospheric fields\n(Z500, T850, wind, humidity)\nfrom ERA5 / NCEP reanalysis", ACCENT),
    ("2. Flow Field", "Learn a velocity field v(x, t)\nthat transforms weather state\nfrom time t to t + dt", ACCENT2),
    ("3. ODE Integration", "Solve dx/dt = v(x, t) using\nadaptive ODE solvers\n(Dormand-Prince, RK4)", ACCENT3),
    ("4. Physics Losses", "Enforce conservation laws:\nmass, energy, enstrophy,\ngeostrophic balance", ACCENT4),
]
for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(3.2)
    add_card(slide, x, Inches(2.2), Inches(2.95), Inches(2.8), title, desc.split("\n"), color)

# Arrow between cards
for i in range(3):
    x = Inches(3.45) + i * Inches(3.2)
    add_text_box(slide, x, Inches(3.2), Inches(0.3), Inches(0.5),
                 ">", font_size=28, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             "Result: Models achieve 99.7% of ECMWF IFS HRES skill on WeatherBench2 benchmarks. "
             "Physics-informed training ensures physically consistent predictions that respect "
             "atmospheric dynamics and conservation laws.",
             font_size=16, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Build System & Developer Guide
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Build System & Developer Guide", font_size=36, bold=True, color=WHITE)

# Makefile targets
add_card(slide, Inches(0.5), Inches(1.3), Inches(4.0), Inches(5.7),
         "Makefile Targets (make ...)",
         ["install           pip install -e .",
          "install-all        Python + frontend + docs",
          "test               Run pytest (100+ tests)",
          "test-frontend      Run Vitest",
          "test-coverage      pytest + coverage report",
          "lint               flake8",
          "format             black + isort",
          "build              sdist + wheel",
          "build-frontend     Vite production build",
          "build-docs         MkDocs build",
          "dev-backend        FastAPI on :8000",
          "dev-frontend       Vite on :5173",
          "dev-streamlit      Streamlit on :8501",
          "docker-build       Build Docker image",
          "release-check      Full pre-release validation",
          "version-check      Verify version consistency",
          "clean              Remove all artifacts"],
         ACCENT)

# Scripts
add_card(slide, Inches(4.8), Inches(1.3), Inches(3.9), Inches(2.7),
         "Standalone Scripts (scripts/)",
         ["build.sh           Build python/frontend/docs",
          "  ./scripts/build.sh python",
          "  ./scripts/build.sh --clean",
          "lint.sh            All linters with --fix mode",
          "  ./scripts/lint.sh --fix",
          "dev.sh             Launch dev servers",
          "  ./scripts/dev.sh all",
          "check_version.py   Version consistency"],
         ACCENT2)

# CI/CD
add_card(slide, Inches(4.8), Inches(4.3), Inches(3.9), Inches(2.7),
         "CI/CD (GitHub Actions)",
         ["tests.yml         pytest + lint on push/PR",
          "publish.yml       Build + publish to PyPI",
          "deploy-pages.yml  Build React + deploy to GH Pages",
          "docs.yml          Build + deploy MkDocs",
          "",
          "All workflows use Makefile targets",
          "All include pip/npm caching"],
         ACCENT3)

# Version files
add_card(slide, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.7),
         "Project Configuration",
         ["Version tracked in 3 files (must match):",
          "  weatherflow/version.py (canonical)",
          "  pyproject.toml",
          "  setup.py",
          "",
          "Build backend: hatchling",
          "Frontend: Vite 5 + React 18",
          "Docs: MkDocs Material",
          "Docker: docker-compose profiles",
          "",
          "Code quality:",
          "  black (formatting)",
          "  isort (import sorting)",
          "  flake8 (linting)",
          "  mypy (type checking)",
          "  pre-commit hooks"],
         ACCENT4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Project Structure
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Project Structure", font_size=36, bold=True, color=WHITE)

tree = [
    ("weatherflow/", "Core Python package", True),
    ("  models/", "Flow matching, physics-guided, stochastic models", False),
    ("  training/", "Trainers, flow trainer, metrics, losses", False),
    ("  data/", "ERA5, NCEP data loaders, streaming, GAIA", False),
    ("  physics/", "Atmospheric constraints, conservation losses", False),
    ("  inference/", "Prediction pipeline, ensemble methods", False),
    ("  server/", "FastAPI REST API (app.py)", False),
    ("  education/", "Graduate atmospheric dynamics tools", False),
    ("  enterprise/", "Custom industry model builder", False),
    ("gcm/", "General Circulation Model", True),
    ("  core/", "Dynamics engine, primitive equations", False),
    ("  physics/", "Radiation, convection, microphysics", False),
    ("applications/", "Domain applications", True),
    ("  renewable_energy/", "Wind + solar power converters", False),
    ("  extreme_event_analysis/", "Heatwave, AR detection", False),
    ("streamlit_app/", "Interactive web app (28 pages)", True),
    ("frontend/", "React + TypeScript SPA", True),
    ("docs/", "MkDocs documentation", True),
    ("tests/", "100+ pytest tests", True),
    ("scripts/", "Build, lint, dev scripts", True),
    ("Makefile", "Unified build system entry point", True),
]

for i, (path, desc, is_header) in enumerate(tree):
    y = Inches(1.3) + i * Inches(0.28)
    color = ACCENT if is_header else LIGHT_GRAY
    fsize = 14 if is_header else 12
    bold = is_header
    tf = add_text_box(slide, Inches(1.0), y, Inches(3.5), Inches(0.3),
                      path, font_size=fsize, bold=bold, color=color,
                      font_name="Consolas")
    add_text_box(slide, Inches(5.0), y, Inches(7.5), Inches(0.3),
                 desc, font_size=fsize-1, color=MID_GRAY if not is_header else LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Workflows
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Key Workflows for New Users", font_size=36, bold=True, color=WHITE)

workflows = [
    ("1. Load Data", [
        "Open Data Manager page in Streamlit",
        "Click 'Download NCEP Sample' (7 MB, ~10 sec)",
        "Data loads automatically into session",
        "Variables: Z500, T850, wind components",
    ], ACCENT),
    ("2. Train a Model", [
        "Go to Training Workflow page",
        "Choose architecture (flow matching recommended)",
        "Set epochs, batch size, learning rate",
        "Click Train - watch live loss curves",
        "Save checkpoint when satisfied",
    ], ACCENT2),
    ("3. Make Predictions", [
        "Open Weather Prediction page",
        "Load your saved checkpoint",
        "Select forecast length (1-14 days)",
        "Generate predictions",
        "View animated weather maps",
    ], ACCENT3),
    ("4. Evaluate & Compare", [
        "Open Model Comparison page",
        "Select models to compare",
        "View WeatherBench2 metrics",
        "Z500 RMSE, ACC, skill scores",
        "Export charts for publications",
    ], ACCENT4),
]

for i, (title, steps, color) in enumerate(workflows):
    x = Inches(0.3) + i * Inches(3.25)
    add_card(slide, x, Inches(1.3), Inches(3.0), Inches(5.5), title, steps, color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Python Module → Streamlit Page Mapping
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Python Module to Streamlit Mapping", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Every Streamlit page executes real Python code. Here is what runs behind each button click:",
             font_size=16, color=LIGHT_GRAY)

mapping = [
    ("Streamlit Page", "Python Module", "What It Does", True),
    ("Training Workflow", "weatherflow.models.flow_matching", "Real PyTorch neural network training", False),
    ("Wind Power", "applications.renewable_energy.wind_power", "Turbine power curve calculations", False),
    ("Solar Power", "applications.renewable_energy.solar_power", "PVlib solar irradiance calculations", False),
    ("Extreme Events", "applications.extreme_event_analysis.detectors", "Scientific event detection algorithms", False),
    ("GCM Simulation", "gcm.core.model.GCM", "Full atmospheric physics simulation", False),
    ("Education", "weatherflow.education.graduate_tool", "Atmospheric dynamics calculations", False),
    ("Flow Matching", "weatherflow.models.flow_matching", "Flow model exploration + ODE solvers", False),
    ("Physics Losses", "weatherflow.physics.losses", "Conservation law enforcement", False),
    ("Hurricane Tracking", "streamlit_app.hurricane_data_utils", "IBTrACS/HURDAT2 + AI inference", False),
    ("Model Library", "weatherflow.model_library", "Pre-trained model zoo management", False),
    ("WeatherBench2", "weatherflow.evaluation", "Standard verification metrics", False),
]

for i, (page, module, action, is_header) in enumerate(mapping):
    y = Inches(1.7) + i * Inches(0.42)
    c1 = WHITE if is_header else LIGHT_GRAY
    c2 = WHITE if is_header else ACCENT
    c3 = WHITE if is_header else MID_GRAY
    fs = 14 if is_header else 13
    bld = is_header
    add_text_box(slide, Inches(0.5), y, Inches(2.8), Inches(0.4),
                 page, font_size=fs, bold=bld, color=c1)
    add_text_box(slide, Inches(3.5), y, Inches(5.3), Inches(0.4),
                 module, font_size=fs, bold=bld, color=c2, font_name="Consolas")
    add_text_box(slide, Inches(9.0), y, Inches(4.0), Inches(0.4),
                 action, font_size=fs, bold=bld, color=c3)

    if is_header:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.5), y + Inches(0.38),
                                       Inches(12.3), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = MID_GRAY
        line.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Getting Help & Resources
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Getting Help & Resources", font_size=36, bold=True, color=WHITE)

add_card(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.5),
         "Documentation",
         ["docs/ - Full MkDocs site with tutorials",
          "docs/tutorials/quickstart.md - 5-min intro",
          "docs/tutorials/era5.md - Working with ERA5",
          "docs/api/ - Complete API reference",
          "make build-docs - Build docs locally"],
         ACCENT)

add_card(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.5),
         "Code & Community",
         ["GitHub: github.com/monksealseal/weatherflow",
          "Issues: Report bugs or request features",
          "License: MIT (free for any use)",
          "Author: Eduardo Siman (monksealseal)",
          "Contact: esiman@msn.com"],
         ACCENT2)

add_card(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.8),
         "Useful Commands",
         ["make help              See all 40+ targets",
          "make test              Run full test suite",
          "make version-check     Verify version consistency",
          "make dev               See dev server options",
          "make release-check     Full pre-release validation",
          "streamlit run streamlit_app/Home.py"],
         ACCENT3)

add_card(slide, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.8),
         "Key Files to Read First",
         ["README.md - Project overview & Streamlit guide",
          "Makefile - Build system reference",
          "weatherflow/version.py - Current version",
          "streamlit_app/Home.py - App entry point",
          "weatherflow/models/flow_matching.py - Core model",
          "tests/ - Examples of how to use the library"],
         ACCENT4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You / Summary
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "WeatherFlow", font_size=54, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
             "Democratizing Weather AI",
             font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

summary_items = [
    "150+ Python modules  |  28 Streamlit pages  |  100+ tests",
    "Real scientific data  |  Physics-informed models  |  No code required",
    "Wind & solar energy  |  Extreme events  |  Hurricane tracking  |  Climate simulation",
]
y = Inches(3.8)
for item in summary_items:
    add_text_box(slide, Inches(1), y, Inches(11), Inches(0.5),
                 item, font_size=17, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    y += Inches(0.5)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "pip install weatherflow   |   streamlit run streamlit_app/Home.py",
             font_size=18, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")

add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
             "github.com/monksealseal/weatherflow   |   MIT License   |   v0.4.3",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────

output_path = "/home/user/weatherflow/WeatherFlow_Onboarding.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
